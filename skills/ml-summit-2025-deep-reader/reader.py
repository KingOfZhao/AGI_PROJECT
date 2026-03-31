#!/usr/bin/env python3
"""
ml-summit-2025-deep-reader — 核心解析引擎
处理 2025MLSUMMITBJ 目录下的 PDF/PPT，生成由浅入深的 Markdown
"""

import re
import json
from pathlib import Path
from typing import List, Dict, Any, Optional
from collections import Counter

import fitz  # PyMuPDF
from pptx import Presentation  # python-pptx


class MLSummitReader:
    """ML Summit 2025 深度阅读器核心引擎"""

    # 技术术语词典（用于智能识别）
    TECH_TERMS = {
        "AI/ML基础": ["大模型", "LLM", "GPT", "Transformer", "attention", "预训练", "微调",
                       "fine-tuning", "RLHF", "alignment", "scaling law", "参数量", "推理",
                       "inference", "训练", "training", "token", "embedding"],
        "Agent": ["agent", "智能体", "MCP", "tool", "function calling", "planning",
                  "多模态", "multimodal", "autonomous", "code agent", "Agentic"],
        "RAG": ["RAG", "检索增强", "retrieval", "知识库", "vector", "embedding",
                "chunking", "rerank"],
        "强化学习": ["强化学习", "RL", "reward", "奖励", "policy", "PPO", "DPO",
                     "RLVR", "reward model", "reinforcement"],
        "推理优化": ["量化", "quantization", "蒸馏", "distillation", "剪枝", "pruning",
                     "KV cache", "投机采样", "speculative", "显存优化", "推理加速",
                     "vLLM", "TensorRT"],
        "计算机视觉": ["视觉", "vision", "图像", "image", "检测", "detection",
                       "OCR", "视频", "video", "3D", "segmentation"],
        "自然语言": ["NLP", "文本", "text", "生成", "generation", "对话", "dialogue",
                     "翻译", "translation", "摘要", "summarization"],
        "自动驾驶": ["自动驾驶", "L4", "感知", "规划", "控制", "激光雷达",
                     "lidar", "端到端", "end-to-end"],
        "基础设施": ["分布式", "distributed", "GPU", "CUDA", "集群", "cluster",
                     "训练框架", "DeepSpeed", "Megatron", "部署", "deployment"],
    }

    def __init__(self, root_dir: str = "/Users/administruter/Downloads/2025MLSUMMITBJ"):
        self.root = Path(root_dir)
        self.output_dir = self.root / "markdown_output"
        self.output_dir.mkdir(exist_ok=True)

    def scan_files(self, pdf_only: bool = False, pptx_only: bool = False) -> List[Path]:
        """扫描所有 PDF 和 PPTX 文件"""
        files = []
        if not pptx_only:
            files.extend(self.root.glob("*.pdf"))
        if not pdf_only:
            files.extend(self.root.glob("*.pptx"))
            files.extend(self.root.glob("*.ppt"))
        return sorted(files)

    def extract_speaker(self, filename: str) -> str:
        """从文件名提取演讲者（格式: 演讲者-标题.pdf）"""
        name = Path(filename).stem
        if "—" in name or "-" in name:
            sep = "—" if "—" in name else "-"
            parts = name.split(sep, 1)
            return parts[0].strip()
        return "未知"

    def extract_title_from_filename(self, filename: str) -> str:
        """从文件名提取标题"""
        name = Path(filename).stem
        for sep in ["—", "-", "_"]:
            if sep in name:
                parts = name.split(sep, 1)
                return parts[-1].strip()
        return name

    # ─── PDF 解析 ───

    def extract_pdf(self, pdf_path: Path) -> Dict[str, Any]:
        """深度解析 PDF"""
        doc = fitz.open(str(pdf_path))
        pages = []
        total_images = 0
        total_chars = 0

        for page_num, page in enumerate(doc):
            text = page.get_text("text").strip()

            # 提取标题（第一行非空文本，去除页码）
            lines = [l.strip() for l in text.split('\n') if l.strip()]
            title = ""
            for line in lines:
                # 跳过纯数字（页码）
                if re.match(r'^\d{1,3}$', line):
                    continue
                title = line
                break

            # 提取图片
            images = len(page.get_images(full=True))
            total_images += images

            # 提取链接
            links = [link["uri"] for link in page.get_links() if link.get("uri")]

            chars = len(text)
            total_chars += chars

            pages.append({
                "page": page_num + 1,
                "title": title,
                "text": text,
                "images": images,
                "links": links,
                "chars": chars,
            })

        doc.close()

        return {
            "type": "pdf",
            "filename": pdf_path.name,
            "speaker": self.extract_speaker(pdf_path.name),
            "topic": self.extract_title_from_filename(pdf_path.name),
            "total_pages": len(pages),
            "total_images": total_images,
            "total_chars": total_chars,
            "avg_density": round(total_chars / max(len(pages), 1), 0),
            "pages": pages,
        }

    # ─── PPTX 解析 ───

    def extract_pptx(self, ppt_path: Path) -> Dict[str, Any]:
        """深度解析 PPTX"""
        prs = Presentation(str(ppt_path))
        slides = []
        total_images = 0
        total_chars = 0

        for slide_num, slide in enumerate(prs.slides):
            texts = []
            images_count = 0
            has_table = False
            has_chart = False

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
                if shape.shape_type == 13:  # Picture
                    images_count += 1
                if shape.has_table:
                    has_table = True
                if shape.has_chart:
                    has_chart = True

            full_text = "\n".join(texts)
            lines = [l.strip() for l in full_text.split('\n') if l.strip()]
            title = lines[0] if lines else ""

            chars = len(full_text)
            total_chars += chars
            total_images += images_count

            slides.append({
                "slide": slide_num + 1,
                "title": title,
                "text": full_text,
                "images": images_count,
                "has_table": has_table,
                "has_chart": has_chart,
                "chars": chars,
            })

        return {
            "type": "pptx",
            "filename": ppt_path.name,
            "speaker": self.extract_speaker(ppt_path.name),
            "topic": self.extract_title_from_filename(ppt_path.name),
            "total_slides": len(slides),
            "total_images": total_images,
            "total_chars": total_chars,
            "avg_density": round(total_chars / max(len(slides), 1), 0),
            "slides": slides,
        }

    # ─── 智能摘要 ───

    def generate_auto_summary(self, data: Dict) -> str:
        """基于前几页+最后页生成智能摘要"""
        units = data.get("pages") or data.get("slides") or []
        if not units:
            return "（内容为空）"

        # 取前3页+最后1页的文本
        key_texts = []
        for u in units[:3]:
            text = u.get("text", "").strip()
            if text:
                key_texts.append(text[:500])
        if len(units) > 3:
            last = units[-1].get("text", "").strip()
            if last:
                key_texts.append(last[:500])

        combined = "\n".join(key_texts)

        # 提取关键句（包含技术关键词的句子）
        sentences = re.split(r'[。\n]', combined)
        tech_sentences = [s.strip() for s in sentences if len(s.strip()) > 10
                         and any(kw in s for kw in ["技术", "模型", "方法", "系统", "框架",
                                                      "训练", "推理", "性能", "结果", "优化",
                                                      "挑战", "创新", "实验", "平台"])]

        # 选取最有信息量的前5句
        selected = sorted(tech_sentences, key=len, reverse=True)[:5]
        if not selected:
            # 回退: 取前3个最长句
            selected = sorted([s.strip() for s in sentences if len(s.strip()) > 15],
                            key=len, reverse=True)[:3]

        return "；".join(selected) + "。" if selected else combined[:200] + "..."

    # ─── 术语提取 ───

    def extract_terms(self, data: Dict) -> List[Dict]:
        """提取关键技术术语及其出现频率"""
        units = data.get("pages") or data.get("slides") or []
        all_text = " ".join(u.get("text", "") for u in units)

        term_counter = Counter()
        term_pages: Dict[str, set] = {}

        unit_key = "page" if data["type"] == "pdf" else "slide"

        for u in units:
            text = u.get("text", "")
            unit_num = u.get(unit_key, 0)
            for category, terms in self.TECH_TERMS.items():
                for term in terms:
                    if term.lower() in text.lower():
                        term_counter[term] += 1
                        if term not in term_pages:
                            term_pages[term] = set()
                        term_pages[term].add(unit_num)

        # 取Top 15
        top_terms = term_counter.most_common(15)
        result = []
        for term, count in top_terms:
            result.append({
                "term": term,
                "count": count,
                "pages": sorted(term_pages.get(term, [])),
            })
        return result

    # ─── 关键点提取 ───

    def extract_key_points(self, text: str) -> List[str]:
        """从文本提取关键点（bullet points）"""
        lines = [l.strip() for l in text.split('\n') if l.strip()]
        key_points = []
        for line in lines:
            # 识别bullet point或编号列表
            if re.match(r'^[•\-\*▸▶]|^\d+[.、）)]', line) and len(line) > 10:
                key_points.append(line)
            elif any(kw in line for kw in ["核心", "关键", "重点", "挑战", "创新", "突破",
                                            "优势", "总结", "结论", "目标", "方法"]):
                if len(line) > 15:
                    key_points.append(line)

        # 去重并限制数量
        seen = set()
        unique = []
        for p in key_points:
            if p not in seen and len(unique) < 5:
                seen.add(p)
                unique.append(p)
        return unique

    # ─── Markdown 生成（4层结构）───

    def build_deep_markdown(self, data: Dict[str, Any]) -> str:
        """生成由浅入深的 Markdown"""
        lines = []

        # 头部
        speaker = data.get("speaker", "未知")
        topic = data.get("topic", "")
        total_units = data.get("total_pages") or data.get("total_slides", 0)
        unit_type = "页" if data["type"] == "pdf" else "幻灯片"

        lines.append(f"# {data['filename']}")
        lines.append(f"\n> 👤 **演讲者**: {speaker} | 📄 {data['type'].upper()} | "
                     f"📊 {total_units} {unit_type} | 🖼️ {data.get('total_images', 0)} 张图片")
        lines.append(f"\n---\n")

        # ═══ Level 1: 浅层摘要 ═══
        lines.append("## 1. 浅层摘要\n")
        lines.append("### 元信息\n")
        lines.append(f"- **文件**: {data['filename']}")
        lines.append(f"- **演讲者**: {speaker}")
        lines.append(f"- **主题**: {topic}")
        lines.append(f"- **类型**: {data['type'].upper()}")
        lines.append(f"- **{unit_type}数**: {total_units}")
        lines.append(f"- **图片数**: {data.get('total_images', 0)}")
        lines.append(f"- **总字符数**: {data.get('total_chars', 0):,}")
        lines.append(f"- **平均密度**: {data.get('avg_density', 0):.0f} 字符/{unit_type[:-1]}\n")

        lines.append("### 智能摘要\n")
        summary = self.generate_auto_summary(data)
        lines.append(f"{summary}\n")

        # ═══ Level 2: 中层关键点 ═══
        units = data.get("pages") or data.get("slides") or []
        unit_key = "page" if data["type"] == "pdf" else "slide"
        unit_label = "页" if data["type"] == "pdf" else "Slide"

        lines.append(f"## 2. 中层关键点（逐{unit_label}）\n")

        for u in units:
            num = u.get(unit_key, 0)
            title = u.get("title", "无标题") or "无标题"
            text = u.get("text", "").strip()

            # 清理文本（去除多余空行）
            text_clean = re.sub(r'\n{3,}', '\n\n', text)

            # 截取预览
            preview = text_clean[:400] + "..." if len(text_clean) > 400 else text_clean

            lines.append(f"### {unit_label} {num} — {title}\n")
            lines.append(f"{preview}\n")

            # 关键点
            key_points = self.extract_key_points(text_clean)
            if key_points:
                lines.append(f"**🔑 关键点:**\n")
                for kp in key_points:
                    lines.append(f"- {kp}")
                lines.append("")

        # ═══ Level 3: 深层技术分析 ═══
        lines.append("## 3. 深层技术分析\n")

        # 术语提取
        terms = self.extract_terms(data)
        if terms:
            lines.append("### 关键术语\n")
            lines.append("| 术语 | 出现次数 | 涉及{0} |\n|------|---------|--------|".format(unit_label))
            for t in terms:
                pages_str = ", ".join(str(p) for p in t["pages"][:8])
                if len(t["pages"]) > 8:
                    pages_str += "..."
                lines.append(f"| {t['term']} | {t['count']} | {pages_str} |")
            lines.append("")

        # 核心论点（基于术语频率最高的主题推断）
        term_categories: Dict[str, int] = {}
        for t in terms:
            for cat, cat_terms in self.TECH_TERMS.items():
                if t["term"] in cat_terms:
                    term_categories[cat] = term_categories.get(cat, 0) + t["count"]

        if term_categories:
            lines.append("### 技术主题分布\n")
            for cat, score in sorted(term_categories.items(), key=lambda x: x[1], reverse=True)[:5]:
                bar = "█" * min(int(score / 2), 20)
                lines.append(f"- **{cat}**: {bar} ({score})\n")

        # ═══ Level 4: 跨文件关联 ═══
        lines.append("## 4. 跨文件关联\n")
        lines.append("> 运行 `python3 main.py --cross-file` 生成跨文件关联报告\n")

        return "\n".join(lines)

    # ─── 单文件处理 ───

    def process_file(self, file_path: Path) -> Path:
        """处理单个文件"""
        print(f"  解析: {file_path.name}")

        if file_path.suffix.lower() == ".pdf":
            data = self.extract_pdf(file_path)
        else:
            data = self.extract_pptx(file_path)

        md = self.build_deep_markdown(data)

        output = self.output_dir / f"{file_path.stem}_deep.md"
        output.write_text(md, encoding="utf-8")

        print(f"  ✅ 生成: {output.name} ({len(md):,} 字符)")
        return output

    # ─── 批量处理 ───

    def process_all(self, pdf_only: bool = False, pptx_only: bool = False):
        """批量处理"""
        files = self.scan_files(pdf_only=pdf_only, pptx_only=pptx_only)
        if not files:
            print(f"❌ 在 {self.root} 未找到文件")
            return

        print(f"\n📂 发现 {len(files)} 个文件")
        print(f"📁 输出目录: {self.output_dir}\n")

        results = []
        for f in files:
            try:
                out = self.process_file(f)
                results.append({"file": f.name, "output": str(out), "status": "ok"})
            except Exception as e:
                print(f"  ❌ 失败: {f.name} — {e}")
                results.append({"file": f.name, "error": str(e), "status": "error"})

        # 汇总
        ok = sum(1 for r in results if r["status"] == "ok")
        fail = sum(1 for r in results if r["status"] == "error")
        print(f"\n{'='*50}")
        print(f"✅ 成功: {ok} | ❌ 失败: {fail} | 📊 总计: {len(results)}")
        print(f"📁 输出: {self.output_dir}")

        return results

    # ─── 跨文件关联 ───

    def generate_cross_file_report(self) -> str:
        """生成跨文件关联报告"""
        files = self.scan_files()
        if not files:
            return "无文件"

        all_data = []
        for f in files:
            try:
                if f.suffix.lower() == ".pdf":
                    data = self.extract_pdf(f)
                else:
                    data = self.extract_pptx(f)
                terms = self.extract_terms(data)
                all_data.append({"file": f.name, "speaker": data.get("speaker", ""),
                                "terms": terms})
            except Exception:
                pass

        # 构建术语→文件矩阵
        term_to_files: Dict[str, List[str]] = {}
        for d in all_data:
            for t in d["terms"]:
                if t["term"] not in term_to_files:
                    term_to_files[t["term"]] = []
                term_to_files[t["term"]].append(f"{d['speaker']}({d['file'][:20]})")

        # 生成报告
        lines = ["# 跨文件关联报告 — 2025 ML Summit\n"]
        lines.append(f"**文件总数**: {len(all_data)}\n")

        # 高频术语（多文件出现）
        multi_file_terms = {t: fs for t, fs in term_to_files.items() if len(fs) >= 2}
        lines.append(f"## 跨文件高频术语（≥2个文件出现）\n")
        lines.append(f"共 {len(multi_file_terms)} 个术语\n")

        sorted_terms = sorted(multi_file_terms.items(), key=lambda x: len(x[1]), reverse=True)
        for term, file_list in sorted_terms[:50]:
            lines.append(f"\n### {term}（{len(file_list)} 个文件）\n")
            for f in file_list[:5]:
                lines.append(f"- {f}")
            if len(file_list) > 5:
                lines.append(f"- ...及其他 {len(file_list)-5} 个文件")

        # 技术主题聚类
        lines.append(f"\n## 技术主题聚类\n")
        topic_clusters: Dict[str, List[str]] = {}
        for term in multi_file_terms:
            for cat, cat_terms in self.TECH_TERMS.items():
                if term in cat_terms:
                    if cat not in topic_clusters:
                        topic_clusters[cat] = []
                    topic_clusters[cat].append(term)

        for cat, terms in sorted(topic_clusters.items(), key=lambda x: len(x[1]), reverse=True):
            lines.append(f"\n### {cat}（{len(terms)} 个跨文件术语）\n")
            for t in terms:
                lines.append(f"- **{t}**: {len(term_to_files[t])} 个文件")

        report = "\n".join(lines)
        report_path = self.output_dir / "cross_file_report.md"
        report_path.write_text(report, encoding="utf-8")
        print(f"📄 跨文件关联报告: {report_path}")

        return report
